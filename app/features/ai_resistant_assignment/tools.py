import os 
from typing import List
import pypdf
import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from langchain_core.prompts import PromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.output_parsers import StrOutputParser

load_dotenv()

current_dir = os.path.dirname(os.path.abspath(__file__))
prompt_path = os.path.join(current_dir, "prompts", "base_prompt.txt")

with open(prompt_path, "r", encoding="utf-8") as f:
    prompt_template = f.read()

def extract_text_from_url(url: str) -> str:
    """Extract text content from a URL"""
    try:
        if "youtube.com" in url or "youtu.be" in url:
            return f"YouTube video content from: {url} (Note: Implement YouTube transcript extraction)"
        
        else:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text[:5000]
            
    except Exception as e:
        return f"Error extracting content from URL: {str(e)}"
    
def validate_input_format(input_path: str) -> bool:
    supported_formats = [".csv", ".pdf", ".docx", ".pptx", ".txt"]
    return input_path.startswith("http") or os.path.splitext(input_path)[1].lower() in supported_formats

def extract_text_from_input(input_path: str) -> str:
    if input_path.startswith("http"):
        return extract_text_from_url(input_path)
    
    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".pdf":
        with open(input_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            return " ".join(page.extract_text() for page in reader.pages)
    elif ext == ".docx":
        doc = Document(input_path)
        return " ".join(p.text for p in doc.paragraphs)
    elif ext == ".pptx":
        prs = Presentation(input_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text_frame = getattr(shape, "text_frame", None)
                        if text_frame and text_frame.text:
                            texts.append(text_frame.text)
        return " ".join(texts)
    elif ext in [".csv", ".txt"]:
        with open(input_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        raise ValueError("Unsupported File Format")
    

def get_chain():
    llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-pro",
        temperature=0.7, 
        max_tokens=600,
        google_api_key=os.getenv("api_key")
    )
    prompt = PromptTemplate.from_template(prompt_template)
    output_parser = StrOutputParser()
    chain = prompt | llm | output_parser
    return chain

def parse_ideas_from_response(response: str, grade_level: str) -> List[dict]:
    ideas = response.split("Update to make this Assignment AI-resistant (idea)")
    results = []
    for i in ideas[1:]:
        text = i.strip()
        if text.startswith(":"):
            text = text[1:].strip()
        if "Explanation:" in text:
            idea, explanation = text.split("Explanation:", 1)
        else:
            idea, explanation = text, "No Explanation provided"
        results.append({
            "assignment_idea": idea.strip(),
            "explanation": explanation.strip()
        })
    return results

